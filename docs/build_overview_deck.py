"""Build docs/AIForgeCrew-Overview.pptx.

The deck is a binary, so it cannot be reviewed in a diff. This script is the
reviewable source of truth for it: change the script, re-run, commit both.

    uv run --with python-pptx python docs/build_overview_deck.py

DESIGN RULE: diagram first, prose last. A slide earns its text only when a
picture cannot carry the point. Every figure is drawn from native PowerPoint
shapes rather than an exported image, so the file stays ~70 KB, the text stays
searchable, and anyone can drag a box in PowerPoint without going back to a
design tool.

Numbers here are measured on main, not estimated; see the memory notes behind
each slide. If a figure changes, change it here and re-run.
"""
from __future__ import annotations

from contextlib import suppress
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "AIForgeCrew-Overview.pptx"

# ---------------------------------------------------------------- palette --
INK = RGBColor(0x1A, 0x1A, 0x1A)      # near-black body text
MUTED = RGBColor(0x6B, 0x6B, 0x6B)    # captions, footers
LINE = RGBColor(0xC9, 0xC9, 0xC9)     # hairlines and connectors
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# One accent family. Boxes are tinted, never saturated, so text stays readable
# in a printed handout and on a projector that crushes contrast.
BLUE = RGBColor(0x1F, 0x4E, 0x79)
BLUE_T = RGBColor(0xDC, 0xE7, 0xF2)
GREEN = RGBColor(0x1E, 0x6B, 0x52)
GREEN_T = RGBColor(0xDC, 0xEE, 0xE7)
AMBER = RGBColor(0x8A, 0x5A, 0x00)
AMBER_T = RGBColor(0xF7, 0xEC, 0xD5)
PLUM = RGBColor(0x5B, 0x3A, 0x6E)
PLUM_T = RGBColor(0xEA, 0xE1, 0xF0)
RED = RGBColor(0x8C, 0x2F, 0x2F)
RED_T = RGBColor(0xF7, 0xDF, 0xDF)
SLATE = RGBColor(0x3A, 0x3A, 0x3A)
SLATE_T = RGBColor(0xEC, 0xEC, 0xEC)

W, H = Inches(13.333), Inches(7.5)     # 16:9

_slide_no = 0


def _e(v) -> Emu:
    """Coerce a coordinate to integer EMU.

    Length arithmetic silently produces floats (`Inches(2.85) / 2`), and
    python-pptx writes those straight into the XML as `1988820.0`, which is not
    a valid ST_Coordinate — PowerPoint rejects the file. Every helper funnels
    its geometry through here so a float can never reach the document.
    """
    return Emu(int(round(float(v))))


# ------------------------------------------------------------- primitives --
def _txbody(tf, size, color, bold, align, space_after=0):
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Segoe UI"
        p.space_after = Pt(space_after)


def text(slide, x, y, w, h, s, size=12, color=INK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """A plain textbox. Returns the shape so callers can tweak it."""
    box = slide.shapes.add_textbox(_e(x), _e(y), _e(w), _e(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = s.split("\n")
    tf.text = lines[0]
    for extra in lines[1:]:
        tf.add_paragraph().text = extra
    _txbody(tf, size, color, bold, align)
    return box


def box(slide, x, y, w, h, label, sub="", fill=SLATE_T, edge=SLATE,
        size=12, sub_size=9.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        label_color=None, bold=True):
    """A tinted rounded box with a bold label and an optional caption under it.

    This is the deck's workhorse: nearly every figure is boxes plus arrows.
    """
    sp = slide.shapes.add_shape(shape, _e(x), _e(y), _e(w), _e(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = edge
    sp.line.width = Pt(1.0)
    sp.shadow.inherit = False
    with suppress(IndexError, KeyError):   # rounded corners: gentler radius
        sp.adjustments[0] = 0.10
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.text = label
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    for r in p0.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = label_color or edge
        r.font.name = "Segoe UI"
    if sub:
        for line in sub.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(sub_size)
                r.font.bold = False
                r.font.color.rgb = INK
                r.font.name = "Segoe UI"
    return sp


def arrow(slide, x1, y1, x2, y2, color=LINE, width=1.5, dashed=False):
    """A straight connector with an arrowhead, drawn as a thin rotated shape."""
    from pptx.oxml.ns import qn
    conn = slide.shapes.add_connector(1, _e(x1), _e(y1), _e(x2), _e(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle",
                                            "w": "med", "len": "med"})
    ln.append(tail)
    if dashed:
        d = ln.makeelement(qn("a:prstDash"), {"val": "sysDash"})
        ln.insert(0, d)
    return conn


def chip(slide, x, y, w, s, fill=WHITE, edge=LINE, color=INK, size=9.5):
    """A small pill for dense label clouds (tool names, integrations)."""
    return box(slide, x, y, w, Inches(0.30), s, fill=fill, edge=edge,
               size=size, label_color=color, bold=False)


def slide_base(prs, title, kicker="", footer_note=""):
    """A blank white slide with the standard title block and page furniture."""
    global _slide_no
    s = prs.slides.add_slide(prs.slide_layouts[6])          # 6 = blank
    _slide_no += 1
    text(s, Inches(0.6), Inches(0.34), Inches(12.2), Inches(0.5),
         title, size=27, bold=True)
    if kicker:
        text(s, Inches(0.6), Inches(0.92), Inches(12.2), Inches(0.32),
             kicker, size=12.5, color=MUTED)
    # hairline under the title block
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.28),
                            Inches(12.13), _e(Pt(0.75)))
    ln.fill.solid()
    ln.fill.fore_color.rgb = LINE
    ln.line.fill.background()
    ln.shadow.inherit = False
    if footer_note:
        text(s, Inches(0.6), Inches(6.92), Inches(10.6), Inches(0.35),
             footer_note, size=9.5, color=MUTED)
    text(s, Inches(12.05), Inches(6.92), Inches(0.7), Inches(0.3),
         str(_slide_no), size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)
    return s


def legend(slide, x, y, items):
    """Row of colour swatches: [(fill, edge, 'meaning'), ...]."""
    cx = x
    for fill, edge, lab in items:
        sw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, _e(cx),
                                    _e(y), Inches(0.17), Inches(0.17))
        sw.fill.solid()
        sw.fill.fore_color.rgb = fill
        sw.line.color.rgb = edge
        sw.shadow.inherit = False
        text(slide, _e(cx + Inches(0.24)), _e(y - Inches(0.02)), Inches(2.3),
             Inches(0.22), lab, size=9.5, color=MUTED)
        cx += Inches(0.30) + Inches(len(lab) * 0.062)


# ------------------------------------------------------------------ deck ---
def title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    text(s, Inches(0.9), Inches(2.45), Inches(11.5), Inches(1.0),
         "AIForgeCrew", size=54, bold=True)
    text(s, Inches(0.95), Inches(3.45), Inches(11.5), Inches(0.5),
         "A self-hosted coding-agent platform — how it is built and how it works",
         size=17, color=MUTED)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(4.12),
                             Inches(1.5), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False
    y = Inches(4.6)
    for i, (lab, sub) in enumerate([
            ("Ticket → PR", "19 agent roles"),
            ("Chat agent", "109 tools, 3 modes"),
            ("Memory", "files, not a DB"),
            ("Fleet sync", "groups + redaction")]):
        box(s, Inches(0.95) + i * Inches(2.75), y, Inches(2.5), Inches(0.92),
            lab, sub, fill=WHITE, edge=LINE, size=13, label_color=BLUE)
    text(s, Inches(0.95), Inches(6.95), Inches(8), Inches(0.3),
         "Self-hosted · no telemetry · no database to operate",
         size=10.5, color=MUTED)


def architecture(prs):
    s = slide_base(
        prs, "Architecture",
        "One FastAPI process on port 8799. Everything else is a library inside it.",


                "State on disk: ~/.aiforge/ — agent_config.json · integrations.json · memory/ · "
                "work/ · "
            "aiforge.db (SQLite)")
    L, WD = Inches(0.75), Inches(11.85)

    box(s, L, Inches(1.55), WD, Inches(0.62), "React UI",
        "Chat · Tickets · Library · Settings · Memory sync",
        fill=SLATE_T, edge=SLATE, size=13)
    arrow(s, Inches(6.7), Inches(2.17), Inches(6.7), Inches(2.42))

    box(s, L, Inches(2.42), WD, Inches(0.62), "FastAPI  ·  port 8799",
        "REST + SSE streaming   ·   approval gate   ·   runtime settings   ·   jobs scheduler",
        fill=BLUE_T, edge=BLUE, size=13)

    # four engines side by side
    cols = [("Chat agent", "ReAct loop\n109 tools", GREEN_T, GREEN),
            ("ADK pipeline", "19 agents\ngit worktrees", PLUM_T, PLUM),
            ("Memory", "OKR-DAG\n+ SQLite", AMBER_T, AMBER),
            ("Integrations", "Jira · GitLab\nConfluence · MCP", SLATE_T, SLATE)]
    cw, gap = Inches(2.85), Inches(0.15)
    for i, (lab, sub, f, e) in enumerate(cols):
        x = L + i * (cw + gap)
        arrow(s, x + cw / 2, Inches(3.04), x + cw / 2, Inches(3.30))
        box(s, x, Inches(3.30), cw, Inches(1.15), lab, sub, fill=f, edge=e,
            size=13)
        arrow(s, x + cw / 2, Inches(4.45), x + cw / 2, Inches(4.72))

    box(s, L, Inches(4.72), WD, Inches(0.60), "LLM client",
        "rate ceiling  ·  request meter  ·  model registry  ·  automatic cloud escalation",
        fill=BLUE_T, edge=BLUE, size=13)
    arrow(s, Inches(6.7), Inches(5.32), Inches(6.7), Inches(5.58))
    box(s, L, Inches(5.58), WD, Inches(0.60), "Any OpenAI-compatible endpoint",
        "LM Studio  ·  vLLM  ·  mlx  ·  OpenRouter  ·  cloud key",
        fill=WHITE, edge=LINE, size=13, label_color=INK)


def integrations(prs):
    s = slide_base(
        prs, "Integrations",
        "One HTTP integration layer. A tool with no credentials is hidden, not broken.",


                "Unconfigured tools return *_not_configured with a hint and drop out of the "
                "catalogue — "
            "a missing credential never crashes a turn.")

    box(s, Inches(5.05), Inches(3.05), Inches(3.2), Inches(1.0),
        "AIForgeCrew", "runtime/tools + integrations.json",
        fill=BLUE_T, edge=BLUE, size=15)

    # Ring of integration groups around the hub.
    ring = [
        ("Jira",        "21 tools\nissues · sprints · JQL", GREEN_T, GREEN, 0.75, 1.62),
        ("Confluence",  "14 tools\npages · attachments",    GREEN_T, GREEN, 3.60, 1.62),
        ("GitLab",      "MRs · pipelines\njobs · files",    GREEN_T, GREEN, 6.45, 1.62),
        ("GitHub",      "pull requests",                    GREEN_T, GREEN, 9.30, 1.62),
        ("Email",       "send · read",                      AMBER_T, AMBER, 0.75, 5.05),
        ("Web",         "search · fetch\ncrawl · browser",  AMBER_T, AMBER, 3.60, 5.05),
        ("MCP",         "any MCP server",                   PLUM_T, PLUM,   6.45, 5.05),
        ("Local host",  "shell · LSP · tests\ntypecheck · IPython", SLATE_T, SLATE, 9.30, 5.05),
    ]
    for lab, sub, f, e, x, y in ring:
        box(s, Inches(x), Inches(y), Inches(2.6), Inches(0.95), lab, sub,
            fill=f, edge=e, size=12.5)
        # connect each box to the hub
        hub_x, hub_y = Inches(6.65), Inches(3.55)
        bx, by = Inches(x + 1.3), Inches(y + (0.95 if y < 3 else 0.0))
        arrow(s, bx, by, hub_x, hub_y if y < 3 else Inches(4.05),
              color=LINE, width=1.25)

    legend(s, Inches(0.75), Inches(6.45), [
        (GREEN_T, GREEN, "work trackers / VCS"),
        (AMBER_T, AMBER, "outside world"),
        (PLUM_T, PLUM, "extensible"),
        (SLATE_T, SLATE, "the machine itself"),
    ])


def chat_modes(prs):
    s = slide_base(
        prs, "How chat works — Simple, Plan, Team",
        "Three modes, one engine. The mode decides who runs and what may be written.",

            "Plain-text ReAct, so no vendor function-calling is required and any model can drive "
            "it. A turn can be stopped and resumed, not restarted.")

    # Shared turn loop across the top.
    steps = ["Build context", "Model call", "Gate", "Run tool", "Observe", "FINAL"]
    subs = ["context_bundle()", "THOUGHT / ACTION", "policy + risk",
            "diff shown", "condense if near limit", "+ next-step chip"]
    x0, bw, gap = Inches(0.75), Inches(1.83), Inches(0.20)
    text(s, Inches(0.75), Inches(1.45), Inches(4), Inches(0.25),
         "THE TURN LOOP  (all three modes)", size=10, bold=True, color=MUTED)
    for i, (st, sb) in enumerate(zip(steps, subs, strict=True)):
        x = x0 + i * (bw + gap)
        box(s, x, Inches(1.75), bw, Inches(0.75), st, sb,
            fill=BLUE_T, edge=BLUE, size=11.5, sub_size=8.5)
        if i < len(steps) - 1:
            arrow(s, x + bw, Inches(2.12), x + bw + gap, Inches(2.12))
    text(s, Inches(0.75), Inches(2.58), Inches(11.9), Inches(0.22),
         "loops until FINAL, the step cap, the turn deadline, or Stop",
         size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # Three mode lanes.
    lanes = [
        ("SIMPLE", "one agent, full tools", GREEN_T, GREEN,
         ["Single ReAct agent", "Reads and WRITES files",
          "Quick mode caps steps", "Approvals: off by default"],
         "Everyday edits, fixes, questions"),
        ("PLAN", "read-only, no writes", AMBER_T, AMBER,
         ["Same loop, write tools removed", "Produces a plan, not a diff",
          "memory_lookup active", "Approvals: ON"],
         "Decide before you spend tokens"),
        ("TEAM", "the full 19-agent pipeline", PLUM_T, PLUM,
         ["ADK pipeline without a ticket", "Fans out to 4 worktrees",
          "One SPEC.md, then reconciled", "Approvals: off by default"],
         "Multi-file builds from one ask"),
    ]
    lw, lgap = Inches(3.83), Inches(0.20)
    for i, (name, tag, f, e, bullets, when) in enumerate(lanes):
        x = x0 + i * (lw + lgap)
        box(s, x, Inches(3.02), lw, Inches(0.62), name, tag,
            fill=f, edge=e, size=15)
        for j, b in enumerate(bullets):
            chip(s, x, Inches(3.78) + j * Inches(0.38), lw, b)
        box(s, x, Inches(5.42), lw, Inches(0.55), when, "",
            fill=WHITE, edge=e, size=10.5, label_color=e, bold=False)

    box(s, x0, Inches(6.15), Inches(11.9), Inches(0.55),
        "Team mode = the ticket pipeline with the ticket removed",

            "Triage → Enhancer → Architect → Planner → Verifier → Doer ↔ Refiner ↔ Feedback → "
            "Validator → Live-verifier → Learner",
        fill=SLATE_T, edge=SLATE, size=11.5, sub_size=9.5)


def memory(prs):
    s = slide_base(
        prs, "Memory — markdown files, not a database",
        "The whole store is a folder you can read, grep, diff and delete.",

            "Recall works with nothing installed: keyword + spell-correction by default; semantic "
            "KNN is one opt-in flag (static embeddings ~30 MB, sqlite-vec, no torch).")

    # write path
    text(s, Inches(0.75), Inches(1.45), Inches(5), Inches(0.25),
         "WRITE PATH", size=10, bold=True, color=MUTED)
    wp = [("capture()", "tagged by topic\n+ which agent wrote it"),
          ("topic brief", "LLM dedupe + merge\nsplit on oversize"),
          ("OKF node", "markdown + frontmatter\ntyped edges")]
    for i, (lab, sub) in enumerate(wp):
        x = Inches(0.75) + i * Inches(2.55)
        box(s, x, Inches(1.75), Inches(2.25), Inches(0.85), lab, sub,
            fill=GREEN_T, edge=GREEN, size=12, sub_size=8.5)
        if i < 2:
            arrow(s, x + Inches(2.25), Inches(2.17), x + Inches(2.55), Inches(2.17))

    # the store
    box(s, Inches(0.75), Inches(3.00), Inches(3.55), Inches(2.55),
        "~/.aiforge/memory/okf/", "", fill=WHITE, edge=AMBER, size=12.5,
        label_color=AMBER)
    for i, d in enumerate(["objectives/   — why",
                           "key_results/  — what",
                           "learnings/    — how",
                           "sessions/     — already ran"]):
        text(s, Inches(1.00), Inches(3.55) + i * Inches(0.42), Inches(3.1),
             Inches(0.3), d, size=11)
    text(s, Inches(1.00), Inches(5.18), Inches(3.1), Inches(0.3),
         "+ aiforge.db  (SQLite: FTS, vectors)", size=10, color=MUTED)

    arrow(s, Inches(4.30), Inches(4.28), Inches(4.72), Inches(4.28))

    box(s, Inches(4.72), Inches(3.00), Inches(3.55), Inches(2.55),
        "In-memory graph", "", fill=WHITE, edge=BLUE, size=12.5,
        label_color=BLUE)
    for i, d in enumerate(["typed frontmatter edges",
                           "plain dicts, built at load",
                           "no DB server, no daemon",
                           "git-diffable, greppable"]):
        text(s, Inches(4.97), Inches(3.55) + i * Inches(0.42), Inches(3.1),
             Inches(0.3), "· " + d, size=11)

    arrow(s, Inches(8.27), Inches(4.28), Inches(8.69), Inches(4.28))

    box(s, Inches(8.69), Inches(3.00), Inches(3.90), Inches(2.55),
        "Surgical retrieval", "", fill=WHITE, edge=PLUM, size=12.5,
        label_color=PLUM)
    for i, d in enumerate(["ascend → the objective (why)",
                           "descend → the key result (what)",
                           "constraints + recent activity",
                           "raw code chunks demoted so",
                           "   curated knowledge outranks RAG"]):
        text(s, Inches(8.94), Inches(3.55) + i * Inches(0.38), Inches(3.5),
             Inches(0.3), d, size=10.5)

    # Sits to the RIGHT of the write path, which ends at x=8.10.
    box(s, Inches(8.40), Inches(1.75), Inches(4.22), Inches(0.85),
        "Why it is simpler",
        
            "no vector-DB service to install, tune or back up\ninspect with cat and grep  ·  "
            "move a machine by copying a folder",
        fill=SLATE_T, edge=SLATE, size=12, sub_size=9)


def central_memory(prs):
    s = slide_base(
        prs, "Central memory and groups",
        "One admin hub serves several independent fleets. Only durable knowledge travels.",

            "Spokes initiate in both directions, so only the admin needs a reachable address — "
            "laptops behind NAT just work. Unset AIFORGE_ADMIN_URL = this box IS the admin.")

    # spokes
    for i, (lab, sub) in enumerate([("Laptop A", "compacts locally"),
                                    ("NUC / server", "compacts locally"),
                                    ("Laptop B", "compacts locally")]):
        box(s, Inches(0.75), Inches(1.70) + i * Inches(1.15), Inches(2.5),
            Inches(0.85), lab, sub, fill=GREEN_T, edge=GREEN, size=12.5)
        arrow(s, Inches(3.25), Inches(2.12) + i * Inches(1.15),
              Inches(3.95), Inches(2.12) + i * Inches(1.15))

    # redaction gate on the client side
    box(s, Inches(3.95), Inches(1.70), Inches(2.25), Inches(3.15),
        "redact\n(on the client)",
        "\nsecrets\nprivate\nnoise\n\nblocks a note WHOLE,\nnever edits it",
        fill=RED_T, edge=RED, size=12, sub_size=9.5)
    arrow(s, Inches(6.20), Inches(3.27), Inches(6.95), Inches(3.27))

    # admin
    box(s, Inches(6.95), Inches(1.70), Inches(2.6), Inches(1.35),
        "Admin hub", "merges only\nnever compacts",
        fill=BLUE_T, edge=BLUE, size=14)
    arrow(s, Inches(8.25), Inches(3.05), Inches(8.25), Inches(3.40))
    box(s, Inches(6.95), Inches(3.40), Inches(2.6), Inches(1.45),
        "groups/<g>/", "peers/  mesh/\nokf/  .snapshots/",
        fill=WHITE, edge=BLUE, size=12.5, label_color=BLUE)

    # group isolation
    box(s, Inches(10.05), Inches(1.70), Inches(2.55), Inches(3.15),
        "Group isolation",
            "\nContextVar in the path root\nrepoints the whole tree\n\nmanifests · merges ·\n"
            "snapshots follow without\nknowing groups exist\n\none fold per group",
        fill=PLUM_T, edge=PLUM, size=12.5, sub_size=9)

    facts = [
        ("Discovered, not configured",
         "client reads the group list from the admin URL; one group = auto-join"),
        ("okf/ is unwritable by sync",


                 "any network-driven write below it raises; destination set is provably {peers/, "
                 "mesh/, "
             "okf/.tomb/}"),
        ("Revert is an endpoint",

             "hardlink snapshot before every fold and every pull — and the revert is itself "
             "revertible"),
        ("Ordering is a per-node counter",

             "not a timestamp: a clock-based last-writer-wins hands every conflict to the worst "
             "clock in the fleet"),
    ]
    for i, (h, d) in enumerate(facts):
        y = Inches(5.10) + i * Inches(0.44)
        text(s, Inches(0.75), y, Inches(3.4), Inches(0.3), "· " + h,
             size=10.5, bold=True)
        text(s, Inches(4.20), y, Inches(8.4), Inches(0.3), d, size=10.5,
             color=MUTED)


def security(prs):
    s = slide_base(
        prs, "Security and supply chain",
        "No telemetry. Nothing leaves the machine unless you point it somewhere.",

            "Stated plainly: by default the chat agent has full, unsandboxed shell access on the "
            "host. The workspace jail guards file tools, not arbitrary shell — run untrusted work "
            "in container mode.")

    # scanner results, as tiles
    text(s, Inches(0.75), Inches(1.45), Inches(6), Inches(0.25),
         "SCANNED AND TRIAGED", size=10, bold=True, color=MUTED)
    tiles = [("0", "Sonar bugs", GREEN_T, GREEN),
             ("0", "Trivy findings", GREEN_T, GREEN),
             ("9", "vulns — all\ndocumented TLS\nopt-outs", AMBER_T, AMBER),
             ("186", "dependency\npackages", SLATE_T, SLATE),
             ("86%", "line coverage\n~4,000 tests", GREEN_T, GREEN),
             ("≤15", "cognitive\ncomplexity, every\nfunction", GREEN_T, GREEN)]
    for i, (big, lab, f, e) in enumerate(tiles):
        x = Inches(0.75) + i * Inches(2.02)
        sp = box(s, x, Inches(1.75), Inches(1.85), Inches(1.15), big, lab,
                 fill=f, edge=e, size=22, sub_size=8.5)
        del sp

    # CI scan pipeline
    text(s, Inches(0.75), Inches(3.15), Inches(6), Inches(0.25),
         "CI SCAN PIPELINE  (.gitlab-ci.yml)", size=10, bold=True, color=MUTED)
    box(s, Inches(0.75), Inches(3.45), Inches(2.3), Inches(0.95), "build",
        "uv sync --frozen\nuv.lock decides versions", fill=BLUE_T, edge=BLUE,
        size=12.5, sub_size=8.5)
    arrow(s, Inches(3.05), Inches(3.92), Inches(3.45), Inches(3.92))
    box(s, Inches(3.45), Inches(3.45), Inches(2.3), Inches(0.95), "sonar/",
        "coverage.xml · junit.xml\n+ web/src included", fill=GREEN_T,
        edge=GREEN, size=12.5, sub_size=8.5)
    box(s, Inches(3.45), Inches(4.55), Inches(2.3), Inches(0.95), "blackduck/",
        "pinned requirements\nCycloneDX SBOM · wheels", fill=PLUM_T,
        edge=PLUM, size=12.5, sub_size=8.5)
    arrow(s, Inches(3.05), Inches(3.92), Inches(3.45), Inches(5.02))
    arrow(s, Inches(5.75), Inches(3.92), Inches(6.15), Inches(3.92))
    box(s, Inches(6.15), Inches(3.45), Inches(2.2), Inches(0.95), "SonarQube",
        "scanner-cli\nGIT_DEPTH 0", fill=WHITE, edge=GREEN, size=12.5,
        sub_size=8.5, label_color=GREEN)
    arrow(s, Inches(5.75), Inches(5.02), Inches(6.15), Inches(5.02))
    box(s, Inches(6.15), Inches(4.55), Inches(2.2), Inches(0.95), "Black Duck",
        "detect10.sh\nships SHIPPED pins", fill=WHITE, edge=PLUM, size=12.5,
        sub_size=8.5, label_color=PLUM)

    # posture column
    box(s, Inches(8.75), Inches(3.45), Inches(3.85), Inches(2.05),
        "Data control", "", fill=WHITE, edge=SLATE, size=12.5,
        label_color=SLATE)
    for i, d in enumerate([
            "Telemetry: none, and none in the lock",
            "Prompts go only to YOUR endpoint",
            "At rest: ~/.aiforge/ — delete it, it is gone",
            "Loopback by default; loud warning if not",
            "SSRF guard on fetch and crawl",
            "Workspace jail on by default"]):
        text(s, Inches(8.95), Inches(3.90) + i * Inches(0.26), Inches(3.55),
             Inches(0.25), "· " + d, size=9.5)

    box(s, Inches(0.75), Inches(5.72), Inches(11.87), Inches(0.60),
        "The 9 vulnerabilities are decisions, not debt",


                "TLS verification is disabled only per endpoint, opt-in, for self-signed internal "
                "hosts "
            "(Jira / Confluence / GitLab) — each one individually triaged and documented",
        fill=AMBER_T, edge=AMBER, size=12, sub_size=9.5)


def ticket_pipeline(prs):
    s = slide_base(
        prs, "How a ticket becomes a PR",
        "The same tools, driven by 19 specialised roles instead of one.",

            "Context gatherers run in parallel — researcher, repo-map (tree-sitter + PageRank over "
            "the AST), conventions.")
    stages = ["Triage", "Enhancer", "Architect", "Planner", "Verifier"]
    x0, bw, gap = Inches(0.75), Inches(2.20), Inches(0.20)
    for i, st in enumerate(stages):
        x = x0 + i * (bw + gap)
        box(s, x, Inches(1.60), bw, Inches(0.72), st, "", fill=BLUE_T,
            edge=BLUE, size=13)
        if i < len(stages) - 1:
            arrow(s, x + bw, Inches(1.96), x + bw + gap, Inches(1.96))

    box(s, x0, Inches(2.75), Inches(11.85), Inches(1.55),
        "Build loop  —  up to 4 parallel subtasks, each in its own git worktree",

            "Doer edits  ↔  tests run  ↔  Feedback reads the failures  ↔  Refiner fixes\n"
            "all building against one SPEC.md written up front, then reconciled",
        fill=GREEN_T, edge=GREEN, size=13, sub_size=10)

    for i, (st, sub) in enumerate([("Validator", "gates the result"),
                                   ("Live-verifier", "runs the real recipe"),
                                   ("Learner", "writes memory back")]):
        x = x0 + i * Inches(4.02)
        arrow(s, x + Inches(1.85), Inches(4.30), x + Inches(1.85), Inches(4.62))
        box(s, x, Inches(4.62), Inches(3.70), Inches(0.85), st, sub,
            fill=PLUM_T, edge=PLUM, size=13)

    box(s, x0, Inches(5.75), Inches(11.85), Inches(0.60), "Pull request",

            "the Learner's objectives, key results and learnings are LLM-verified before they are "
            "saved",
        fill=SLATE_T, edge=SLATE, size=13, sub_size=9.5)


def deploy(prs):
    s = slide_base(
        prs, "How to deploy it",
        "Zero-Docker is a first-class mode, not a degraded one.",
        "One service, one port, one folder of state. Moving a machine is copying ~/.aiforge/.")
    modes = [
        ("Laptop", "./run.sh", "everything local\nunset admin URL = you are the admin",
         GREEN_T, GREEN),
        ("NUC / server", "systemd user unit", "always-on admin hub\nspokes sync to it",
         BLUE_T, BLUE),
        ("Container", "docker compose", "for untrusted work\nthe jail you actually want",
         PLUM_T, PLUM),
    ]
    for i, (lab, cmd, sub, f, e) in enumerate(modes):
        x = Inches(0.75) + i * Inches(4.02)
        box(s, x, Inches(1.70), Inches(3.70), Inches(1.9), lab,
            f"\n{cmd}\n\n{sub}", fill=f, edge=e, size=17, sub_size=10.5)

    box(s, Inches(0.75), Inches(3.95), Inches(11.85), Inches(0.62),
        "Requires", 
            "Python 3.12  ·  one OpenAI-compatible endpoint  ·  "
            "no GPU, no torch, no Postgres, no Neo4j",
        fill=WHITE, edge=LINE, size=13, sub_size=10.5, label_color=INK)

    box(s, Inches(0.75), Inches(4.80), Inches(11.85), Inches(1.35),
        "Kept honest", 
            "~4,000 tests, 0 failing  ·  86% line coverage  ·  "

                "Sonar 0 bugs, every function under cognitive complexity 15\n"
                "chat and pipeline tool registries are parity-tested  ·  "
            "uv.lock pinned and shipped  ·  CycloneDX SBOM in CI",
        fill=SLATE_T, edge=SLATE, size=13, sub_size=10.5)


def what_it_is(prs):
    s = slide_base(prs, "What it is",
                   "One service. Two ways in. No database to operate.")
    cards = [
        ("Ticket → PR", "a plain-language ticket\nbecomes a pull request",
         "19 specialised agent roles", BLUE_T, BLUE),
        ("Chat agent", "the whole filesystem,\nthree modes",
         "109 tools, plain-text ReAct", GREEN_T, GREEN),
        ("Any endpoint", "OpenAI-compatible\nis the only requirement",
         "LM Studio · vLLM · cloud", AMBER_T, AMBER),
        ("Memory is files", "markdown + one SQLite,\nnot a service",
         "greppable, git-diffable", PLUM_T, PLUM),
    ]
    for i, (lab, sub, foot, f, e) in enumerate(cards):
        x = Inches(0.75) + i * Inches(3.03)
        box(s, x, Inches(1.70), Inches(2.80), Inches(2.05), lab,
            f"\n{sub}", fill=f, edge=e, size=16, sub_size=11)
        text(s, x, Inches(3.85), Inches(2.80), Inches(0.3), foot,
             size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

    box(s, Inches(0.75), Inches(4.45), Inches(11.87), Inches(0.95),
        "Runs on a laptop, a NUC, or a container",
        
            "zero-Docker mode is a first-class deployment, not a degraded one  ·  "
            "no GPU, no torch, no Postgres, no Neo4j",
        fill=SLATE_T, edge=SLATE, size=14, sub_size=10.5)

    nums = [("465", "Python modules"), ("~102k", "lines in core"),
            ("109", "chat tools"), ("19", "agent roles"),
            ("~4,000", "tests, 0 failing"), ("86%", "line coverage"),
            ("186", "packages"), ("0", "databases")]
    for i, (big, lab) in enumerate(nums):
        x = Inches(0.75) + i * Inches(1.49)
        text(s, x, Inches(5.62), Inches(1.35), Inches(0.4), big,
             size=20, bold=True, align=PP_ALIGN.CENTER, color=BLUE)
        text(s, x, Inches(6.05), Inches(1.35), Inches(0.3), lab,
             size=9, color=MUTED, align=PP_ALIGN.CENTER)
    text(s, Inches(0.75), Inches(6.45), Inches(11.87), Inches(0.25),
         "measured on main, not estimated", size=9, color=MUTED,
         align=PP_ALIGN.CENTER)


def tool_surface(prs):
    s = slide_base(
        prs, "The tool surface",
        "109 tools in the chat registry; per-agent allowlists in the pipeline.",


                "Two registries on purpose — the chat agent and the pipeline Doer are separate "
                "engines; "
            "drift between them fails a parity test and a startup check.")
    groups = [
        ("Files", GREEN_T, GREEN,
         ["file_read", "file_write", "file_patch", "editor", "multi_edit", "list_dir"]),
        ("Code", BLUE_T, BLUE,
         ["grep", "find", "lsp", "run_command", "run_tests", "typecheck",
          "format", "ipython", "serve"]),
        ("VCS", AMBER_T, AMBER, ["git (targeted)", "github_pr", "gitlab MRs"]),
        ("Integrations", PLUM_T, PLUM,
         ["jira_* (21)", "confluence_* (14)", "gitlab_*", "email",
          "web_search", "crawl", "browser", "MCP"]),
        ("Memory", AMBER_T, AMBER,
         ["memory_lookup", "memory_write", "remember_rule", "skill search"]),
        ("Long-running", SLATE_T, SLATE,
         ["watch_until", "schedule_task", "serve", "plan_progress"]),
    ]
    y = Inches(1.60)
    for name, f, e, tools in groups:
        box(s, Inches(0.75), y, Inches(1.95), Inches(0.62), name, "",
            fill=f, edge=e, size=12.5)
        cx = Inches(2.85)
        for t in tools:
            w = Inches(0.42 + len(t) * 0.085)
            chip(s, cx, y + Inches(0.16), w, t, edge=e, color=e)
            cx += w + Inches(0.12)
        y += Inches(0.80)

    box(s, Inches(0.75), Inches(6.05), Inches(11.87), Inches(0.62),
        "Unconfigured is not broken",
        


                    "a tool with no credentials returns *_not_configured plus a hint and is hidden "
                    "from "
                "the catalogue  ·  "
            "`aiforge-tool <name> '<json>'` gives scripts the same registry, read-only by default",
        fill=WHITE, edge=LINE, size=12, sub_size=9.5, label_color=INK)


def built_on(prs):
    s = slide_base(
        prs, "What it is built on",
        "Borrow the engine, own the judgement. 186 packages, uv.lock pinned and shipped.",
        "Dropped on purpose: pymongo (imported by nothing) and aider-chat "
        "(its RepoMap vendored instead) — 212 -> 186 packages, and no pin pressure.")

    layers = [
        ("SERVICE", BLUE_T, BLUE,
         ["FastAPI + uvicorn", "pydantic", "httpx", "croniter (jobs)"]),
        ("MODEL ACCESS", GREEN_T, GREEN,
         ["openai SDK", "litellm", "any OpenAI-compatible", "NO vendor tool-calling"]),
        ("PIPELINE", PLUM_T, PLUM,
         ["google-adk 2.1.x", "SequentialAgent", "LoopAgent", "git worktrees"]),
        ("MEMORY + CODE", AMBER_T, AMBER,
         ["SQLite + sqlite-vec", "tree-sitter", "vendored RepoMap", "model2vec embeddings"]),
    ]
    cw, gap = Inches(2.85), Inches(0.15)
    for i, (name, f, e, items) in enumerate(layers):
        x = Inches(0.75) + i * (cw + gap)
        box(s, x, Inches(1.55), cw, Inches(0.58), name, "", fill=f, edge=e,
            size=13)
        for j, it in enumerate(items):
            chip(s, x, Inches(2.25) + j * Inches(0.38), cw, it, edge=e,
                 color=INK, size=9)

    # the two things that are ours, not borrowed
    box(s, Inches(0.75), Inches(3.95), Inches(6.05), Inches(1.35),
        "The ReAct loop is ours",
        "plain text THOUGHT / ACTION, not vendor function-calling —\n"
        "so a 4-bit local model drives the same 109 tools a frontier model does",
        fill=WHITE, edge=GREEN, size=14, sub_size=10.5, label_color=GREEN)
    box(s, Inches(7.05), Inches(3.95), Inches(5.57), Inches(1.35),
        "ADK drives the 19-role pipeline",
        "SequentialAgent[ Planner, Verifier, LoopAgent[Doer, Feedback], Learner ]\n"
        "pinned <2.2 — 2.3 added a Workflow validation that rejects the v6 graph",
        fill=WHITE, edge=PLUM, size=14, sub_size=10.5, label_color=PLUM)

    box(s, Inches(0.75), Inches(5.50), Inches(11.87), Inches(0.95),
        "How a library gets in",
        "it earns its place or it goes: sqlite-vec is CORE, not an extra, because without it every "
        "memory write hit \"no such module: vec0\"  ·  aider-chat was dropped and its RepoMap "
        "vendored under Apache-2.0, proven by a differential render of 7,640 identical bytes",
        fill=SLATE_T, edge=SLATE, size=13, sub_size=9.5)


def message_flow(prs):
    s = slide_base(
        prs, "What happens when you type a message",
        "Two passes, opposite directions — then the loop.",
        "Most messages skip the save step entirely: a cheap keyword pre-check decides, "
        "so an ordinary turn costs no extra model call.")

    box(s, Inches(0.75), Inches(1.55), Inches(1.95), Inches(0.70), "message",
        "", fill=WHITE, edge=LINE, size=13, label_color=INK)
    arrow(s, Inches(2.70), Inches(1.90), Inches(3.05), Inches(1.90))

    # 1 — the write pass
    box(s, Inches(3.05), Inches(1.42), Inches(9.57), Inches(0.96),
        "1.  SAVE   —   anything here worth remembering?",
        "pre-filter (free)  ->  classify  ->  rule | memory | feedback | none  +  scope"
        "  ->  store\n"
        "pure directive with no task?  ->  \"Got it, saved.\"  and the agent never runs",
        fill=AMBER_T, edge=AMBER, size=13, sub_size=9.5)
    arrow(s, Inches(7.8), Inches(2.38), Inches(7.8), Inches(2.62))

    # 2 — the read pass
    box(s, Inches(3.05), Inches(2.62), Inches(9.57), Inches(0.96),
        "2.  GATHER   —   what do I already know that applies?",
        "preferences  ·  rules (glob-matched)  ·  skills + workflows (relevance)  ·  "
        "repo map  ·  memory recall\nassembled in ONE place, so adding a context source "
        "touches one function",
        fill=GREEN_T, edge=GREEN, size=13, sub_size=9.5)
    arrow(s, Inches(7.8), Inches(3.58), Inches(7.8), Inches(3.82))

    # 3 — the loop
    box(s, Inches(3.05), Inches(3.82), Inches(9.57), Inches(0.62),
        "3.  THE LOOP", "", fill=BLUE_T, edge=BLUE, size=13)
    steps = [("think", ""), ("pick a tool", ""), ("safety gate", "approve / reject"),
             ("run it", "diff shown"), ("observe", "condense if full")]
    bw = Inches(1.82)
    for i, (lab, sub) in enumerate(steps):
        x = Inches(3.05) + i * (bw + Inches(0.11))
        box(s, x, Inches(4.55), bw, Inches(0.68), lab, sub, fill=WHITE,
            edge=BLUE, size=11, sub_size=8, label_color=BLUE)
        if i < len(steps) - 1:
            arrow(s, x + bw, Inches(4.89), x + bw + Inches(0.11), Inches(4.89))
    text(s, Inches(3.05), Inches(5.30), Inches(9.57), Inches(0.24),
         "repeats until FINAL, the step cap, the turn deadline, or Stop  —  "
         "a stopped turn RESUMES, it does not restart",
         size=9.5, color=MUTED, align=PP_ALIGN.CENTER)

    arrow(s, Inches(7.8), Inches(5.58), Inches(7.8), Inches(5.82))
    box(s, Inches(3.05), Inches(5.82), Inches(9.57), Inches(0.70),
        "4.  REPLY   +   predicted next step",
        "it ACTS on that prediction only inside the blast-radius table — otherwise it offers",
        fill=SLATE_T, edge=SLATE, size=13, sub_size=9.5)


def versus(prs):
    s = slide_base(
        prs, "Where it differs from the other agent tools",
        "It adopts their formats rather than replacing them — then adds the two "
        "things none of them do.",
        "Named comparisons are the ones the codebase itself cites: the SKILL.md gaps vs "
        "Hermes Agent / OpenClaw, and Cursor's glob-scoped rule format.")

    box(s, Inches(0.75), Inches(1.52), Inches(11.87), Inches(0.78),
        "Reads what your repo already has",
        ".cursor/rules/*.mdc  ·  .cursorrules  ·  AGENTS.md  ·  SKILL.md "
        "(agentskills.io)  —  no migration, no lock-in",
        fill=GREEN_T, edge=GREEN, size=14, sub_size=10.5)

    rows = [
        ("Ticket → PR AND chat", "usually one or the other",
         "one service does both, on the same tool registry"),
        ("Model requirement", "vendor function-calling",
         "plain-text ReAct — a 4-bit local model drives it"),
        ("Memory store", "a vector DB to install, tune, back up",
         "markdown + one SQLite file; grep it, delete it"),
        ("Skill discovery", "trigger substring match",
         "searched by RELEVANCE, so the right one is found"),
        ("Gets better how", "you write the prompts",
         "write_skill() / write_workflow() — it authors its own"),
        ("Across machines", "per-machine, or a vendor cloud",
         "fleet sync: groups, redaction, snapshot + revert"),
        ("Telemetry", "usually on by default", "none, and none in the lock"),
    ]
    for i, (k, a, b) in enumerate(rows):
        y = Inches(2.48) + i * Inches(0.58)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), y,
                                Inches(11.87), Inches(0.52))
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE if i % 2 else RGBColor(0xF7, 0xF7, 0xF7)
        bg.line.fill.background()
        bg.shadow.inherit = False
        text(s, Inches(0.90), y + Inches(0.14), Inches(2.6), Inches(0.3), k,
             size=11, bold=True)
        text(s, Inches(3.55), y + Inches(0.14), Inches(4.0), Inches(0.3), a,
             size=10.5, color=MUTED)
        text(s, Inches(7.70), y + Inches(0.14), Inches(4.8), Inches(0.3),
             "▸ " + b, size=10.5, color=INK)
    text(s, Inches(3.55), Inches(2.20), Inches(4.0), Inches(0.28),
         "typical agent tool", size=11, bold=True, color=MUTED)
    text(s, Inches(7.70), Inches(2.20), Inches(4.8), Inches(0.28),
         "AIForgeCrew", size=11, bold=True, color=BLUE)


def rules_skills_workflows(prs):
    s = slide_base(
        prs, "Rules, skills and workflows",
        "How the agent is taught — and how it teaches itself.",
        "All three are markdown on disk, repo-local overriding global by name. "
        "Ships with 3 rules, 4 skills and 1 workflow; the rest you write, or it writes.")

    cols = [
        ("RULES", "glob-scoped, deterministic\nzero LLM cost", BLUE_T, BLUE,
         [".aiforge/rules/*.md",
          ".cursor/rules/*.mdc",
          ".cursorrules  (always)",
          "AGENTS.md  (always)"],
         "matched against the ticket's scope globs — in BOTH directions, so\n"
         "src/** matches src/a/** — then rendered as one capped block into\n"
         "the planner and doer prompts, replacing a paid LLM gathering call"),
        ("SKILLS", "SKILL.md\nagentskills.io convention", GREEN_T, GREEN,
         ["~/.aiforge/skills/<name>/",
          "<repo>/.aiforge/skills/",
          "<repo>/.claude/skills/",
          "frontmatter + markdown body"],
         "searched by RELEVANCE — description and trigger overlap, not a\n"
         "substring match — so the right skill is discovered rather than\n"
         "named. write_skill() authors one the moment something hard is solved"),
        ("WORKFLOWS", "WORKFLOW.md\nsame machinery, own folder", PLUM_T, PLUM,
         ["~/.aiforge/workflows/<name>/",
          "<repo>/.aiforge/workflows/",
          "<repo>/.claude/workflows/",
          "end-to-end procedures"],
         "skills are small reusable how-tos; workflows are the longer recipe —\n"
         "\"how we cut a release\", \"how we triage a flaky test\". Added by\n"
         "write_workflow(), or by dropping a file in by hand"),
    ]
    cw, gap = Inches(3.83), Inches(0.20)
    for i, (name, tag, f, e, roots, note) in enumerate(cols):
        x = Inches(0.75) + i * (cw + gap)
        box(s, x, Inches(1.55), cw, Inches(0.68), name, tag,
            fill=f, edge=e, size=15, sub_size=9)
        for j, r in enumerate(roots):
            chip(s, x, Inches(2.35) + j * Inches(0.36), cw, r, edge=e,
                 color=INK, size=9)
        text(s, x + Inches(0.05), Inches(3.85), cw - Inches(0.10),
             Inches(1.0), note, size=9, color=MUTED)

    # the always-on capture loop that feeds them
    text(s, Inches(0.75), Inches(4.92), Inches(6), Inches(0.22),
         "ALWAYS-ON CAPTURE  —  runs BEFORE the agent, on every chat message",
         size=10, bold=True, color=MUTED)
    steps = [
        ("user message", ""),
        ("classify", "ONE capped LLM pass"),
        ("rule | memory\nfeedback | none", "+ scope, confidence"),
        ("route + store", "repo rules · md_store\nmemory · in-session"),
    ]
    bw = Inches(2.70)
    for i, (lab, sub) in enumerate(steps):
        x = Inches(0.75) + i * (bw + Inches(0.28))
        box(s, x, Inches(5.22), bw, Inches(0.78), lab, sub,
            fill=AMBER_T if i else WHITE, edge=AMBER if i else LINE,
            size=11.5, sub_size=8.5, label_color=None if i else INK)
        if i < 3:
            arrow(s, x + bw, Inches(5.61), x + bw + Inches(0.28), Inches(5.61))

    box(s, Inches(0.75), Inches(6.15), Inches(11.87), Inches(0.55),
        "A gate is NEVER disabled by the classifier",
        "a directive stated in passing is captured deterministically instead of "
        "hoping the model calls remember_rule — but disabling a commit or delete "
        "gate stays a separate, scoped, revocable opt-in",
        fill=RED_T, edge=RED, size=11.5, sub_size=9.5)


def safety(prs):
    s = slide_base(
        prs, "Safety rails",
        "Autonomy is a table over blast radius, not a confidence threshold.",

            "A high-confidence prediction never buys its way up a tier: tier 3 never acts, at any "
            "confidence.")
    tiers = [
        ("TIER 1", "reversible, local", "acts on its own",
         "read a file · grep · run tests · typecheck", GREEN_T, GREEN),
        ("TIER 2", "writes inside the workspace", "acts, shows the diff",
         "edit a file · run a command · git add -p", AMBER_T, AMBER),
        ("TIER 3", "outside, or hard to undo", "OFFERS — never acts",
         "push · deploy · delete · post to Jira / Confluence / email",
         RED_T, RED),
    ]
    for i, (t, radius, behaviour, ex, f, e) in enumerate(tiers):
        y = Inches(1.62) + i * Inches(1.32)
        box(s, Inches(0.75), y, Inches(1.75), Inches(1.08), t, radius,
            fill=f, edge=e, size=15, sub_size=9)
        box(s, Inches(2.70), y, Inches(3.05), Inches(1.08), behaviour, "",
            fill=WHITE, edge=e, size=12.5, label_color=e)
        box(s, Inches(5.95), y, Inches(6.67), Inches(1.08), "", ex,
            fill=WHITE, edge=LINE, size=1, sub_size=11)

    rails = [
        ("Every write shows a diff", "risky commands pause for Approve / Reject"),
        ("`git add -A` is refused", "blanket staging is never what was meant"),
        ("DELETE has its own floor", "the delete guard fires even with approvals off"),
        ("Workspace jail on by default", "AIFORGE_WORKSPACE_DIR clamps file and exec"),
    ]
    for i, (h, d) in enumerate(rails):
        y = Inches(5.65) + (i // 2) * Inches(0.42)
        x = Inches(0.75) + (i % 2) * Inches(6.1)
        text(s, x, y, Inches(2.9), Inches(0.3), "· " + h, size=10.5, bold=True)
        text(s, x + Inches(2.95), y, Inches(3.1), Inches(0.3), d, size=10.5,
             color=MUTED)


def limits(prs):
    s = slide_base(
        prs, "Context, cost and long-running work",
        "The three things that break an agent in production, each with a rail.",

            "0 = no limit is the shipped default for the budgets; the rate ceiling is always on "
            "because a 429 is not a budget question.")
    cols = [
        ("Context", BLUE_T, BLUE,
         ["condense when near the limit",
          "one assembly seam: context_bundle()",
          "compaction runs once a day, 18:00",
          "Cave mode trims for small models"]),
        ("Cost and rate", AMBER_T, AMBER,
         ["sliding-window rpm ceiling",
          "monotonic clock, never starts full",
          "request meter in the toolbar",
          "structured path obeys it too"]),
        ("Long-running", GREEN_T, GREEN,
         ["a turn extends its own runaway cap",
          "progress = new reads + edits + tree change",
          "watch_until · schedule_task · serve",
          "a stopped turn resumes, not restarts"]),
    ]
    for i, (name, f, e, items) in enumerate(cols):
        x = Inches(0.75) + i * Inches(4.02)
        box(s, x, Inches(1.62), Inches(3.70), Inches(0.65), name, "",
            fill=f, edge=e, size=15)
        for j, it in enumerate(items):
            chip(s, x, Inches(2.42) + j * Inches(0.46), Inches(3.70), it,
                 edge=e, color=INK)

    box(s, Inches(0.75), Inches(4.55), Inches(11.87), Inches(0.85),
        "The trap each rail exists for",

            "a token-bucket that starts full delivers 2N in the first minute  ·  a boot-time fold "
            "leaked compaction into the working day  ·  \"distinct tool args\" counted a retry as "
            "progress",
        fill=SLATE_T, edge=SLATE, size=12.5, sub_size=10)

    box(s, Inches(0.75), Inches(5.60), Inches(11.87), Inches(0.95),
        "Escalation",
            "the LLM client carries a model registry and escalates to a cloud model only when the "
            "local one has actually failed the step — not on a timer, and never silently",
        fill=WHITE, edge=BLUE, size=13, sub_size=10.5, label_color=BLUE)


def comparison(prs):
    s = slide_base(
        prs, "Compared with the usual approach",
        "Every operational burden a vector-database stack adds, removed.")
    rows = [
        ("Storage", "vector DB service to install, tune, back up",
         "markdown files + one SQLite under ~/.aiforge/"),
        ("Inspect", "query the DB, decode embeddings", "cat · grep · an editor"),
        ("Version", "opaque rows", "git-diffable; hardlink snapshots + revert endpoint"),
        ("Move a machine", "dump and restore", "copy a folder"),
        ("Delete", "hope the rows are gone", "delete the folder — there is no second copy"),
        ("Model", "vendor function-calling required", "plain-text ReAct — any model drives it"),
        ("Telemetry", "usually on by default", "none, and none in the lock"),
    ]
    hdr_y = Inches(1.55)
    for lab, x, w, col in [("", Inches(0.75), Inches(2.2), MUTED),
                           ("Typical agent stack", Inches(3.05), Inches(4.5), MUTED),
                           ("AIForgeCrew", Inches(7.75), Inches(4.87), BLUE)]:
        if lab:
            text(s, x, hdr_y, w, Inches(0.3), lab, size=11.5, bold=True,
                 color=col)
    for i, (k, a, b) in enumerate(rows):
        y = Inches(1.95) + i * Inches(0.66)
        f = WHITE if i % 2 else RGBColor(0xF7, 0xF7, 0xF7)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), y,
                                Inches(11.87), Inches(0.60))
        bg.fill.solid()
        bg.fill.fore_color.rgb = f
        bg.line.fill.background()
        bg.shadow.inherit = False
        text(s, Inches(0.90), y + Inches(0.17), Inches(2.1), Inches(0.3), k,
             size=11, bold=True)
        text(s, Inches(3.05), y + Inches(0.17), Inches(4.5), Inches(0.3), a,
             size=10.5, color=MUTED)
        text(s, Inches(7.75), y + Inches(0.17), Inches(4.7), Inches(0.3), b,
             size=10.5, color=INK)


def summary(prs):
    s = slide_base(prs, "In one slide", "")
    points = [
        ("One process", "FastAPI on 8799; every engine is a library inside it",
         BLUE_T, BLUE),
        ("Two ways in", "a ticket becomes a PR, or you chat with the filesystem",
         GREEN_T, GREEN),
        ("Three chat modes", "Simple writes, Plan is read-only, Team is the pipeline",
         GREEN_T, GREEN),
        ("Memory is a folder", "markdown + one SQLite; grep it, diff it, delete it",
         AMBER_T, AMBER),
        ("One admin merges", "spokes compact locally; redaction runs before a note is advertised",
         PLUM_T, PLUM),
        ("Nothing phones home", "no telemetry; prompts go only to the endpoint you set",
         RED_T, RED),
    ]
    for i, (h, d, f, e) in enumerate(points):
        x = Inches(0.75) + (i % 2) * Inches(6.05)
        y = Inches(1.70) + (i // 2) * Inches(1.55)
        box(s, x, y, Inches(5.82), Inches(1.25), h, f"\n{d}",
            fill=f, edge=e, size=15, sub_size=11)
    text(s, Inches(0.75), Inches(6.42), Inches(11.87), Inches(0.3),


                 "Self-hosted · no database to operate · ~4,000 tests, 0 failing · 86% coverage · "
                 "Sonar "
             "0 bugs",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    title_slide(prs)
    what_it_is(prs)
    architecture(prs)
    built_on(prs)
    integrations(prs)
    chat_modes(prs)
    message_flow(prs)
    ticket_pipeline(prs)
    tool_surface(prs)
    rules_skills_workflows(prs)
    memory(prs)
    central_memory(prs)
    limits(prs)
    safety(prs)
    security(prs)
    versus(prs)
    comparison(prs)
    deploy(prs)
    summary(prs)
    prs.save(OUT)
    print(f"wrote {OUT}  ({OUT.stat().st_size / 1024:.0f} KB, "
          f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
